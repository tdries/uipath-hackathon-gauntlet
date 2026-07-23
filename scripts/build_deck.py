"""Populate the UiPath AgentHack submission deck with Gauntlet content.

Reads `Submission deck.template.pptx.bak` (the original UiPath master template
with placeholder slides) and writes `Submission deck.pptx`, preserving the
template's master, layouts, and brand styling.

Run with the project venv activated:
    python scripts/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

BLACK = RGBColor(0x17, 0x21, 0x25)  # UiPath ink-dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


REPO = Path(__file__).resolve().parent.parent
# Prefer the original UiPath master template if it's present locally; otherwise
# fall back to the branded base snapshot (a read-only copy of a previously built
# deck, so the master, layouts, and brand styling survive even when the .bak is
# not on this machine).
_BAK = REPO / "Submission deck.template.pptx.bak"
TEMPLATE = _BAK if _BAK.exists() else REPO / "Submission deck.base.pptx"
OUTPUT = REPO / "Submission deck.pptx"
SCREENSHOTS = REPO / "docs" / "screenshots"


# ──────────────────────────── helpers ────────────────────────────


def set_lines(text_frame, lines, font_size_pt=None, color=None, bold=None):
    """Replace text_frame contents with `lines` (one per paragraph). Some
    placeholder cascades break when we rewrite the runs, so optionally pin an
    explicit color/size/bold to make the text render reliably."""
    if not lines:
        return
    text_frame.text = lines[0]
    for line in lines[1:]:
        p = text_frame.add_paragraph()
        p.text = line
    for p in text_frame.paragraphs:
        for r in p.runs:
            if font_size_pt is not None:
                r.font.size = Pt(font_size_pt)
            if color is not None:
                r.font.color.rgb = color
            if bold is not None:
                r.font.bold = bold


def set_cell(cell, text, font_size_pt=None, color=None, bold=None):
    cell.text = text
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            if font_size_pt is not None:
                r.font.size = Pt(font_size_pt)
            if color is not None:
                r.font.color.rgb = color
            if bold is not None:
                r.font.bold = bold


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def find_shape_by_text(slide, needle):
    for shape in slide.shapes:
        if shape.has_text_frame and needle in shape.text_frame.text:
            return shape
    return None


def slide_has_picture(slide):
    return any(getattr(sh, "shape_type", None) == 13 for sh in slide.shapes)


# ──────────────────────────── content ────────────────────────────


def build():
    prs = Presentation(str(TEMPLATE))

    # ─── Slide 1: Title ──────────────────────────────────────────────
    s1 = prs.slides[0]
    eyebrow = find_shape_by_text(s1, "UiPath AgentHack")
    if eyebrow:
        set_lines(eyebrow.text_frame, ["UiPath AgentHack 2026  •  Track 3 (Test Cloud)  •  Finalist"])
    # .bak uses the placeholder "Presentation title"; the branded base already
    # reads "Gauntlet ...". Match either.
    title = find_shape_by_text(s1, "Presentation title") or find_shape_by_text(s1, "Gauntlet")
    if title:
        # Keep the subtitle to one line so it doesn't wrap into the eyebrow below.
        set_lines(title.text_frame, [
            "Gauntlet",
            "Adversarial Test Cloud.",
        ])

    # ─── Slide 2: Team ───────────────────────────────────────────────
    s2 = prs.slides[1]
    picture_shapes = [sh for sh in s2.shapes if sh.shape_type == 13]
    name_shapes = [sh for sh in s2.shapes if sh.has_text_frame and sh.text_frame.text.strip() == "Jane Doe"]
    title_shapes = [sh for sh in s2.shapes if sh.has_text_frame and "Job title @Company" in sh.text_frame.text]

    if name_shapes:
        set_lines(name_shapes[0].text_frame, ["Tim Dries"], font_size_pt=22, color=BLACK, bold=True)
    if title_shapes:
        set_lines(title_shapes[0].text_frame, [
            "Solution Architect @ Biztory",
            "tim.dries@biztory.be",
        ], font_size_pt=14, color=BLACK)

    for pic in picture_shapes[1:]:
        remove_shape(pic)
    for nshape in name_shapes[1:]:
        remove_shape(nshape)
    for tshape in title_shapes[1:]:
        remove_shape(tshape)

    team_footer = find_shape_by_text(s2, "Team name")
    if team_footer:
        set_lines(team_footer.text_frame, ["Biztory  •  Gauntlet"])

    # ─── Slide 3: Problem + Solution ─────────────────────────────────
    s3 = prs.slides[2]
    head = (find_shape_by_text(s3, "Problem statement and proposed solution")
            or find_shape_by_text(s3, "The gap between"))
    if head:
        set_lines(head.text_frame,
                  ["Your agent passes the tests you wrote. Then a real attacker shows up."])

    problem_body = (find_shape_by_text(s3, "What real-world problem")
                    or find_shape_by_text(s3, "UiPath Agent Evaluations test what you wrote"))
    if problem_body:
        set_lines(problem_body.text_frame, [
            "UiPath Agent Evaluations score your agent against a static rubric a human wrote.",
            "They cannot tell you what happens when an attacker invents a prompt nobody imagined.",
            "",
            "Agent safety today is hope, not proof: no adversary in the loop, and no coverage a compliance officer can actually read.",
        ])

    solution_body = (find_shape_by_text(s3, "Brief summary of the solution")
                     or find_shape_by_text(s3, "A Red Coach"))
    if solution_body:
        set_lines(solution_body.text_frame, [
            "Gauntlet puts a Red Coach, a frontier LLM, in the ring against your agent.",
            "When the known attacks stop landing, the Coach invents new ones mid-fight.",
            "Every win becomes a UiPath Test Manager regression test, automatically.",
            "Every loss opens an Action Center task with the fix already drafted.",
            "",
            "Every fight double-tagged to OWASP LLM Top-10 and MITRE ATLAS.",
        ])

    # ─── Slide 4: Benefits and UiPath components ─────────────────────
    s4 = prs.slides[3]

    title4 = (find_shape_by_text(s4, "Benefits and technologies used")
              or find_shape_by_text(s4, "Benefits, impact, and UiPath"))
    if title4:
        set_lines(title4.text_frame, ["What it proved, and what it runs on"])

    sub4 = (find_shape_by_text(s4, "Benefits, impact and outcomes")
            or find_shape_by_text(s4, "Outcomes"))
    if sub4:
        set_lines(sub4.text_frame, ["Outcomes"], font_size_pt=24, color=BLACK, bold=True)

    body4 = (find_shape_by_text(s4, "What does this agent actually achieve")
             or find_shape_by_text(s4, "75+ adversarial fights")
             or find_shape_by_text(s4, "adversarial fights"))
    if body4:
        set_lines(body4.text_frame, [
            "Strict posture held the line: 27 of 27 attacks blocked, average defense 90.6 / 100.",
            "Weaken the policy and it leaks: worst-case 22% attacker win rate, measured not guessed.",
            "Headline breach: a planted \"KYC subsystem\" note talked a naive agent into disclosing an $84,320 balance with zero 2FA.",
            "75 fights run, 9 critical findings auto-diagnosed with a patch proposed.",
            "A regression suite that grows itself, plus an auditor-ready OWASP / MITRE matrix.",
        ], font_size_pt=15, color=BLACK)

    details = (find_shape_by_text(s4, "Details") or find_shape_by_text(s4, "Components"))
    if details:
        set_lines(details.text_frame, ["Components"], font_size_pt=24, color=BLACK, bold=True)

    table_shape = None
    for sh in s4.shapes:
        if sh.has_table:
            table_shape = sh
            break
    if table_shape is not None:
        rows = list(table_shape.table.rows)
        set_cell(rows[0].cells[0], "End user", font_size_pt=14)
        set_cell(rows[0].cells[1], "AI engineers and QA leads shipping UiPath agents", font_size_pt=12)
        set_cell(rows[1].cells[0], "Department", font_size_pt=14)
        set_cell(rows[1].cells[1], "AI / Automation CoE; Security & Compliance", font_size_pt=12)
        set_cell(rows[2].cells[0], "Industries", font_size_pt=14)
        set_cell(rows[2].cells[1], "Banking, insurance, healthcare, regulated SaaS", font_size_pt=12)
        set_cell(rows[3].cells[0], "UiPath products", font_size_pt=14)
        set_cell(rows[3].cells[1],
                 "Test Cloud + Test Manager, Maestro (Case + Flow), Agent Builder, Coded Agents, Coded Apps, Action Center, Automation Cloud, uip CLI",
                 font_size_pt=11)
        set_cell(rows[4].cells[0], "Other tech", font_size_pt=14)
        set_cell(rows[4].cells[1],
                 "Frontier LLM, LangGraph, React 19 + Vite, @uipath/uipath-typescript SDK, OWASP LLM Top-10, MITRE ATLAS, built with an agentic coding tool",
                 font_size_pt=11)

    # ─── Slide 5: Architecture ───────────────────────────────────────
    s5 = prs.slides[4]
    head5 = (find_shape_by_text(s5, "Solution architecture")
             or find_shape_by_text(s5, "How it works"))
    if head5:
        set_lines(head5.text_frame, ["How it works"])

    body5 = (find_shape_by_text(s5, "This slide is optional")
             or find_shape_by_text(s5, "FightArena is a long-lived"))
    if body5:
        # Shrink to left half so the screenshot can sit on the right.
        body5.left = Inches(0.4)
        body5.top = Inches(1.7)
        body5.width = Inches(5.6)
        body5.height = Inches(5.2)
        set_lines(body5.text_frame, [
            "FightArena is a long-lived Maestro Case, one per target agent.",
            "",
            "Each round is a RoundOrchestrator Flow:",
            "Red attacks, Blue responds, the Referee agent scores it.",
            "",
            "The Coach runs a self-play loop:",
            "play the highest expected-reward attack,",
            "or ask the LLM to invent a brand-new persona.",
            "",
            "Red win goes to Test Manager as a regression.",
            "Blue loss goes to the Fix Recommender, then Action Center.",
        ], font_size_pt=15, color=BLACK)

    coach_img = SCREENSHOTS / "02-coach-lab.png"
    if coach_img.exists() and not slide_has_picture(s5):
        slide_w = prs.slide_width
        img_w = Inches(6.4)
        img_left = slide_w - img_w - Inches(0.3)
        img_top = Inches(1.65)
        s5.shapes.add_picture(str(coach_img), img_left, img_top, width=img_w)

    # ─── Slide 6: Coverage + demo ────────────────────────────────────
    s6 = prs.slides[5]
    head6 = (find_shape_by_text(s6, "Miscellaneous")
             or find_shape_by_text(s6, "Coverage, audit"))
    if head6:
        set_lines(head6.text_frame, ["Coverage, audit, and the live demo"])

    body6 = (find_shape_by_text(s6, "If you need extra slides")
             or find_shape_by_text(s6, "Every fight tagged against OWASP"))
    if body6:
        body6.left = Inches(0.4)
        body6.top = Inches(1.7)
        body6.width = Inches(5.6)
        body6.height = Inches(5.2)
        set_lines(body6.text_frame, [
            "Every fight is tagged to OWASP LLM Top-10 and MITRE ATLAS.",
            "The Audit view renders it as a live coverage matrix, red where Blue is weak.",
            "",
            "In Coach Lab the LLM call streams from your browser with your own key. The attack persona is authored live on stage, not pre-recorded.",
            "",
            "Framework-neutral: coded agents, Agent Builder low-code, and an external LangGraph target all fight in the same arena.",
            "",
            "▶  youtu.be/1q9W5SC_fxA",
            "📦  github.com/tdries/uipath-hackathon-gauntlet",
        ], font_size_pt=15, color=BLACK)

    audit_img = SCREENSHOTS / "06-audit.png"
    if audit_img.exists() and not slide_has_picture(s6):
        slide_w = prs.slide_width
        img_w = Inches(6.4)
        img_left = slide_w - img_w - Inches(0.3)
        img_top = Inches(1.65)
        s6.shapes.add_picture(str(audit_img), img_left, img_top, width=img_w)

    # ─── Slide 7: Closing ────────────────────────────────────────────
    s7 = prs.slides[6]
    closing = (find_shape_by_text(s7, "Closing message")
               or find_shape_by_text(s7, "Go safe"))
    if closing:
        set_lines(closing.text_frame, ["Go safe, or go home."])

    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
