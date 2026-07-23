"""Build the Gauntlet finalist pitch deck on the real UiPath master template.

Reads `Submission deck.base.pptx` (UiPath master + 22 branded layouts), clears
its slides, and draws 30 native, editable slides using the UiPath brand palette
pulled straight from the template theme. Content is sourced from README.md and
docs/ARCHITECTURE.md; nothing is invented.

    python scripts/build_pitch_deck.py
    -> "Gauntlet Pitch Deck.pptx"
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

try:
    from PIL import Image as _PILImage
except Exception:  # pragma: no cover
    _PILImage = None

REPO = Path(__file__).resolve().parent.parent
TEMPLATE = REPO / "Submission deck.base.pptx"
OUTPUT = REPO / "Gauntlet Pitch Deck.pptx"
OUTPUT_MD = REPO / "docs" / "PITCH-TALK-TRACK.md"
SHOTS = REPO / "docs" / "screenshots"
LOGO = REPO / "gauntlet-logo.png"

# ── UiPath brand palette (extracted from the template theme) ──────────────
INK = RGBColor(0x17, 0x20, 0x24)      # near-black ink
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9D, 0x9D, 0x9D)     # secondary text
MIST = RGBColor(0xF6, 0xF6, 0xF6)     # light card fill
CLOUD = RGBColor(0xEC, 0xEE, 0xF0)    # slightly darker light fill
ORANGE = RGBColor(0xFA, 0x46, 0x16)   # UiPath signature orange (attacker/red)
RUST = RGBColor(0xA3, 0x22, 0x00)     # dark red
CORAL = RGBColor(0xFA, 0x76, 0x78)    # light red
TEAL = RGBColor(0x1E, 0x64, 0x82)     # defender/blue
CYAN = RGBColor(0x0B, 0xA2, 0xB3)     # heal/fix
LIGHTCYAN = RGBColor(0x5B, 0xCB, 0xDE)
PURPLE = RGBColor(0x8B, 0x28, 0x8A)   # referee/judge
INKSOFT = RGBColor(0x3C, 0x45, 0x4A)  # softened ink for body copy

ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}

EMU_IN = 914400

# ─────────────────────────────── helpers ─────────────────────────────────


def _layouts(prs):
    return {lay.name: lay for lay in prs.slide_masters[0].slide_layouts}


def _clear_slides(prs):
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rid = sid.get(qn("r:id"))
        try:
            prs.part.drop_rel(rid)
        except KeyError:
            pass
        lst.remove(sid)


def _noshadow(shape):
    shape.shadow.inherit = False


def _norm_paragraphs(content):
    """content -> list of paragraphs; each paragraph -> list of run dicts."""
    if isinstance(content, str):
        content = content.split("\n")
    paras = []
    for para in content:
        if isinstance(para, str):
            paras.append([{"t": para}])
        elif isinstance(para, dict):
            paras.append([para])
        else:  # list of run dicts
            paras.append(list(para))
    return paras


def text(slide, l, t, w, h, content, size=16, color=INK, bold=False, italic=False,
         align="l", anchor="t", font=None, spacing=1.06, before=0, after=6, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = ANCH[anchor]
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, runs in enumerate(_norm_paragraphs(content)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align if isinstance(align, PP_ALIGN) else ALIGN[align]
        p.line_spacing = spacing
        if before:
            p.space_before = Pt(before)
        p.space_after = Pt(after)
        for rd in runs:
            r = p.add_run()
            r.text = rd["t"]
            f = r.font
            f.size = Pt(rd.get("size", size))
            f.bold = rd.get("bold", bold)
            f.italic = rd.get("italic", italic)
            f.color.rgb = rd.get("color", color)
            fam = rd.get("font", font)
            if fam:
                f.name = fam
    return box


def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, radius=None,
         shape=MSO_SHAPE.RECTANGLE, rot=None):
    if radius is not None:
        shape = MSO_SHAPE.ROUNDED_RECTANGLE
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    _noshadow(sp)
    if rot is not None:
        sp.rotation = rot
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.text_frame.paragraphs[0].text = ""
    return sp


def card(slide, l, t, w, h, fill=MIST, line=None, line_w=1.0, radius=0.045):
    return rect(slide, l, t, w, h, fill=fill, line=line, line_w=line_w, radius=radius)


def node(slide, l, t, w, h, title, sub=None, fill=INK, tcolor=WHITE, scolor=None,
         radius=0.06, tsize=14, ssize=10.5, line=None, line_w=1.25, anchor="m"):
    """A labelled diagram node (rounded rect with centred title + optional sub)."""
    rect(slide, l, t, w, h, fill=fill, line=line, line_w=line_w, radius=radius)
    body = [{"t": title, "size": tsize, "bold": True, "color": tcolor}]
    content = [body]
    if sub:
        content.append([{"t": sub, "size": ssize, "bold": False,
                         "color": scolor or tcolor}])
    text(slide, l + 0.08, t, w - 0.16, h, content, align="c", anchor=anchor,
         spacing=1.02, after=2)


def chip(slide, l, t, w, label, fill, color=WHITE, size=10.5, h=0.34, bold=True):
    rect(slide, l, t, w, h, fill=fill, radius=0.5)
    text(slide, l, t, w, h, label, size=size, color=color, bold=bold,
         align="c", anchor="m", after=0)


def _line_el(conn, head=None, tail=None, dash=None):
    ln = conn.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    for tag, val in (("a:headEnd", head), ("a:tailEnd", tail)):
        if val:
            ln.append(ln.makeelement(qn(tag), {"type": val, "w": "med", "len": "med"}))


def arrow(slide, x1, y1, x2, y2, color=INK, w=2.0, tail="triangle", head=None, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    _noshadow(c)
    _line_el(c, head=head, tail=tail, dash=dash)
    return c


def kicker(slide, act, label):
    """Orange rule under the title + left-aligned act/section eyebrow beside it,
    kept clear of the master's UiPath logo at top-right."""
    rect(slide, 0.42, 1.26, 0.6, 0.07, fill=ORANGE)
    text(slide, 1.15, 1.11, 8.0, 0.32,
         [[{"t": act, "color": ORANGE, "bold": True, "size": 12},
           {"t": f"    {label.upper()}", "color": GRAY, "bold": True, "size": 12}]],
         align="l", anchor="m", after=0)


def content_slide(deck, title_text, act=None, label=None):
    s = deck.add("Headline only")
    deck.set_title(s, title_text)
    if act:
        kicker(s, act, label or "")
    return s


def kpi(slide, l, t, w, number, label, accent=ORANGE, nsize=40, lsize=11.5, h=1.5):
    card(slide, l, t, w, h, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.05)
    rect(slide, l, t, 0.09, h, fill=accent)  # left accent spine
    text(slide, l + 0.28, t + 0.16, w - 0.4, 0.85, number, size=nsize, bold=True,
         color=accent, anchor="m", after=0)
    text(slide, l + 0.3, t + h - 0.56, w - 0.5, 0.5, label, size=lsize, color=INKSOFT,
         anchor="t", spacing=1.0, after=0)


def footer(slide, note):
    text(slide, 0.42, 7.05, 9.0, 0.32, note, size=9, color=GRAY, after=0)


# ─────────────────────────── deck wrapper ────────────────────────────────


class Deck:
    def __init__(self):
        self.prs = Presentation(str(TEMPLATE))
        _clear_slides(self.prs)
        self.L = _layouts(self.prs)
        self.W = Emu(self.prs.slide_width).inches
        self.H = Emu(self.prs.slide_height).inches

    def add(self, layout_name):
        return self.prs.slides.add_slide(self.L[layout_name])

    def set_title(self, slide, txt):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = txt
                return ph
        return None

    def save(self):
        self.prs.save(str(OUTPUT))


# ═══════════════════════════════ slides ══════════════════════════════════


def s01_title(d):
    s = d.add("Headline only")
    d.set_title(s, "")  # suppress placeholder prompt
    # full-bleed dark hero
    rect(s, -0.1, -0.1, d.W + 0.2, d.H + 0.2, fill=INK)
    rect(s, 0, 0, d.W, 0.16, fill=ORANGE)  # top orange rule
    text(s, 0.9, 1.35, 11.5, 0.5,
         [[{"t": "UIPATH AGENTHACK 2026", "color": WHITE, "bold": True, "size": 13},
           {"t": "     TRACK 3 · TEST CLOUD     FINALIST", "color": ORANGE,
            "bold": True, "size": 13}]], after=0)
    text(s, 0.86, 2.05, 11.6, 1.9, "Gauntlet", size=104, bold=True, color=WHITE,
         after=0, spacing=0.9)
    rect(s, 0.95, 4.02, 2.7, 0.09, fill=ORANGE)
    text(s, 0.9, 4.25, 11.5, 0.9, "Adversarial Test Cloud", size=34, color=WHITE,
         bold=False, after=0)
    text(s, 0.9, 5.25, 11.0, 0.6,
         "Put an adversary in the loop. Turn every attack that lands into a permanent test.",
         size=16.5, color=RGBColor(0xC9, 0xCE, 0xD2), after=0)
    text(s, 0.9, 6.7, 11.6, 0.4,
         [[{"t": "Tim Dries · Biztory", "color": WHITE, "size": 12.5, "bold": True},
           {"t": "     go safe, or go home.", "color": GRAY, "size": 12.5}]], after=0)


def s02_hook(d):
    s = d.add("Headline only")
    d.set_title(s, "")
    rect(s, -0.1, -0.1, d.W + 0.2, d.H + 0.2, fill=MIST)
    rect(s, 0, 0, 0.28, d.H, fill=ORANGE)
    text(s, 1.1, 1.5, 11.0, 0.5, "THE UNCOMFORTABLE TRUTH", size=13, bold=True,
         color=ORANGE, after=0)
    text(s, 1.05, 2.25, 11.2, 2.6,
         [[{"t": "Your agent passes every test ", "color": INK, "size": 40, "bold": True}],
          [{"t": "you ", "color": INK, "size": 40, "bold": True},
           {"t": "wrote.", "color": INK, "size": 40, "bold": True}],
          [{"t": "Then a real attacker shows up.", "color": ORANGE, "size": 40,
            "bold": True}]], spacing=1.05, after=4)
    text(s, 1.1, 5.35, 10.8, 0.8,
         "Static rubrics only measure the failures you already imagined. Attackers "
         "specialise in the ones you did not.", size=17, color=INKSOFT, after=0)


def s03_stakes(d):
    s = content_slide(d, "Agents raise the stakes of a wrong answer", "01", "Why")
    text(s, 0.42, 1.5, 12.4, 0.6,
         "A chatbot that says something wrong is embarrassing. An agent that does "
         "something wrong is a breach.", size=16.5, color=INKSOFT, after=0)
    cols = [
        ("They act", "Agents call tools: move money, transfer calls, read and write "
         "customer records. A bad decision has a blast radius.", ORANGE),
        ("Language is the attack surface", "There is no input to sanitise. The prompt "
         "is the exploit, and it can arrive through data, not just the user.", TEAL),
        ("Nondeterministic", "The same guardrail holds four times and fails the fifth. "
         "One passing eval run proves almost nothing.", PURPLE),
    ]
    x, w, gap = 0.42, 3.99, 0.19
    for i, (h, b, c) in enumerate(cols):
        lx = x + i * (w + gap)
        card(s, lx, 2.35, w, 3.9, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.04)
        rect(s, lx, 2.35, w, 0.11, fill=c)
        rect(s, lx + 0.3, 2.72, 0.5, 0.5, fill=c, radius=0.5)
        text(s, lx + 0.3, 2.74, 0.5, 0.5, str(i + 1), size=20, bold=True, color=WHITE,
             align="c", anchor="m", after=0)
        text(s, lx + 0.3, 3.5, w - 0.6, 0.9, h, size=18, bold=True, color=INK, after=0)
        text(s, lx + 0.3, 4.5, w - 0.6, 1.6, b, size=13, color=INKSOFT, spacing=1.12,
             after=0)
    footer(s, "The cost of a jailbreak is no longer reputational. It is operational.")


def s04_gap(d):
    s = content_slide(d, "Today, agent safety is hope, not proof", "02", "The gap")
    text(s, 0.42, 1.52, 12.4, 0.55,
         [[{"t": "UiPath Agent Evaluations answer one question well: ", "size": 16.5,
            "color": INKSOFT},
           {"t": "“does my agent pass the rubric I wrote?”", "size": 16.5,
            "color": INK, "bold": True}]], after=0)
    # the pipeline you have today
    y = 2.5
    steps = [("Write a\nrubric", GRAY), ("Run the\neval", GRAY),
             ("It\npasses", TEAL), ("Ship it", INK)]
    bx, bw, bh, bgap = 0.42, 1.62, 1.0, 0.42
    for i, (lab, c) in enumerate(steps):
        lx = bx + i * (bw + bgap)
        node(s, lx, y, bw, bh, lab, fill=WHITE, tcolor=c, line=CLOUD, line_w=1.4,
             tsize=13)
        if i < len(steps) - 1:
            arrow(s, lx + bw + 0.05, y + bh / 2, lx + bw + bgap - 0.05, y + bh / 2,
                  color=GRAY, w=2.0)
    # the adversary that never entered the loop
    gx = bx + 4 * bw + 3 * bgap + 0.32
    rect(s, gx, y - 0.12, 0.03, bh + 0.24, fill=CLOUD)
    node(s, gx + 0.42, y, 3.55, bh, "A real attacker",
         sub="never in this loop", fill=RUST, tcolor=WHITE, scolor=CORAL, tsize=15,
         ssize=11)
    text(s, gx + 0.42, y + bh + 0.16, 3.55, 0.5, "?!", size=26, bold=True, color=RUST,
         align="c", after=0)
    # the three holes
    text(s, 0.42, 4.35, 12.4, 0.4, "What that leaves unanswered", size=15, bold=True,
         color=INK, after=0)
    holes = [
        "No adversary is ever in the loop. The test set is a mirror of your own blind spots.",
        "“Did we test prompt injection?” resolves to “maybe”, not a coverage number.",
        "A bug in production means manual repro, then a manually written test, if anyone remembers.",
    ]
    for i, hh in enumerate(holes):
        hy = 4.85 + i * 0.62
        rect(s, 0.55, hy + 0.06, 0.16, 0.16, fill=ORANGE, radius=0.5)
        text(s, 0.85, hy, 11.9, 0.55, hh, size=13.5, color=INKSOFT, after=0)


def s04b_judge(d):
    s = content_slide(d, "The LLM judge is not the hard part", "02",
                      "Beyond the judge")
    text(s, 0.42, 1.52, 12.4, 0.75,
         "UiPath Agent Evaluations already use an LLM as a judge to grade your agent "
         "against a rubric. Gauntlet uses one too, as its Referee. The judge was never "
         "the missing piece.", size=15.5, color=INKSOFT, spacing=1.12, after=0)
    # ── left: what an LLM judge does today ──
    card(s, 0.42, 2.62, 5.6, 3.35, fill=RGBColor(0xE9, 0xF1, 0xF4), radius=0.04)
    text(s, 0.72, 2.82, 5.0, 0.35, "TODAY IN TEST CLOUD", size=12, bold=True,
         color=TEAL, after=0)
    node(s, 0.72, 3.45, 1.56, 0.92, "Inputs you\nwrote", fill=WHITE, tcolor=INK,
         line=CLOUD, line_w=1.4, tsize=12)
    node(s, 2.52, 3.45, 1.48, 0.92, "LLM judge", fill=PURPLE, tcolor=WHITE, tsize=13)
    node(s, 4.24, 3.45, 1.56, 0.92, "Pass /\nfail", fill=WHITE, tcolor=TEAL,
         line=CLOUD, line_w=1.4, tsize=13)
    arrow(s, 2.28, 3.91, 2.52, 3.91, color=INK, w=2.0)
    arrow(s, 4.0, 3.91, 4.24, 3.91, color=INK, w=2.0)
    text(s, 0.72, 4.7, 5.1, 0.9,
         "It grades what you already thought to test. Nothing generates the attack you "
         "did not imagine, and the score goes nowhere.", size=12.5, color=INKSOFT,
         spacing=1.14, after=0)
    # ── right: where Gauntlet goes beyond ──
    card(s, 6.3, 2.62, 6.6, 3.35, fill=RGBColor(0xFF, 0xEC, 0xE6), line=ORANGE,
         line_w=1.25, radius=0.04)
    text(s, 6.62, 2.82, 6.0, 0.35, "WHERE GAUNTLET GOES BEYOND", size=12, bold=True,
         color=RUST, after=0)
    pts = [
        ("An adversary invents the inputs", "A judge can only grade attacks someone "
         "ran. The Red Coach generates novel ones the rubric never imagined."),
        ("Self play escalates the exam", "Pass, and it authors a harder attack. The "
         "bar rises by itself, instead of staying frozen at what you wrote."),
        ("A closed loop, not a score", "Verdicts become Test Manager regressions and "
         "Action Center fixes, so the result does work."),
    ]
    for i, (h, b) in enumerate(pts):
        yy = 3.32 + i * 0.88
        rect(s, 6.62, yy + 0.02, 0.4, 0.4, fill=ORANGE, radius=0.5)
        text(s, 6.62, yy + 0.03, 0.4, 0.4, str(i + 1), size=15, bold=True, color=WHITE,
             align="c", anchor="m", after=0)
        text(s, 7.18, yy - 0.04, 5.5, 0.4, h, size=13.5, bold=True, color=INK, after=0)
        text(s, 7.18, yy + 0.32, 5.55, 0.55, b, size=11.5, color=INKSOFT, spacing=1.08,
             after=0)
    # ── the line ──
    card(s, 0.42, 6.05, 12.5, 0.72, fill=INK, radius=0.06)
    text(s, 0.42, 6.05, 12.5, 0.72,
         [[{"t": "An LLM judge tells you whether you passed. ", "size": 14,
            "color": WHITE},
           {"t": "Gauntlet decides what the exam is, and rewrites it every time you "
            "pass.", "size": 14, "bold": True, "color": ORANGE}]], align="c",
         anchor="m", after=0)


def s05_delta(d):
    s = content_slide(d, "The delta: from static rubric to live adversary", "02",
                      "The gap")
    rows = [
        ("Test design", "Static rubrics a human wrote once",
         "Self play that invents new evals every run"),
        ("Coverage", "“Did we test prompt injection?” → maybe",
         "OWASP + MITRE coverage matrix on every release"),
        ("Prod bug", "Manual repro, then a manual test case",
         "Auto populated Test Manager regression"),
        ("The fix", "A developer patches the prompt by hand",
         "Fix Recommender drafts the patch + Action Center task"),
    ]
    top, lw, mw, rw, rh, gap = 2.05, 2.15, 4.55, 5.2, 0.98, 0.14
    # headers
    text(s, 0.42 + lw + 0.1, top - 0.5, mw, 0.4, "TODAY", size=13, bold=True,
         color=GRAY, after=0)
    text(s, 0.42 + lw + 0.1 + mw + 0.2, top - 0.5, rw, 0.4, "WITH GAUNTLET", size=13,
         bold=True, color=ORANGE, after=0)
    for i, (k, a, b) in enumerate(rows):
        y = top + i * (rh + gap)
        text(s, 0.42, y, lw, rh, k, size=13.5, bold=True, color=INK, anchor="m",
             after=0)
        card(s, 0.42 + lw + 0.1, y, mw, rh, fill=MIST, radius=0.05)
        text(s, 0.42 + lw + 0.3, y, mw - 0.4, rh, a, size=13, color=INKSOFT,
             anchor="m", after=0)
        bx = 0.42 + lw + 0.1 + mw + 0.2
        card(s, bx, y, rw, rh, fill=RGBColor(0xFF, 0xEC, 0xE6), line=ORANGE,
             line_w=1.25, radius=0.05)
        text(s, bx + 0.2, y, rw - 0.4, rh,
             [[{"t": b, "size": 13, "color": INK, "bold": True}]], anchor="m", after=0)


def s06_section_meet(d):
    s = d.add("Section Header")
    rect(s, -0.1, -0.1, d.W + 0.2, d.H + 0.2, fill=INK)
    rect(s, 0.7, 2.35, 1.9, 0.11, fill=ORANGE)
    text(s, 0.68, 2.7, 11.0, 0.55, "SO WE FIGHT FIRE WITH FIRE", size=15, bold=True,
         color=ORANGE, after=0)
    text(s, 0.66, 3.3, 12.0, 1.7, "Meet Gauntlet", size=66, bold=True, color=WHITE,
         after=0)
    text(s, 0.7, 5.0, 11.2, 0.8,
         "An adversarial arena where an AI red team attacks your agent until it "
         "breaks, and every break becomes a test.", size=18,
         color=RGBColor(0xC9, 0xCE, 0xD2), after=0)


def s07_what(d):
    s = content_slide(d, "One loop: attack, judge, learn, harden", "03", "What it is")
    y = 2.5
    bw, bh = 2.05, 1.15
    mid = y + bh / 2
    xs = [0.42, 2.86, 5.30]
    node(s, xs[0], y, bw, bh, "Red Coach", sub="AI adversary", fill=ORANGE,
         tcolor=WHITE, scolor=RGBColor(0xFF, 0xD9, 0xCC), tsize=16, ssize=10.5)
    node(s, xs[1], y, bw, bh, "Blue target", sub="your agent", fill=TEAL, tcolor=WHITE,
         scolor=LIGHTCYAN, tsize=16, ssize=10.5)
    node(s, xs[2], y, bw, bh, "Referee", sub="scores the round", fill=PURPLE,
         tcolor=WHITE, scolor=RGBColor(0xE7, 0xC4, 0xE6), tsize=16, ssize=10.5)
    arrow(s, xs[0] + bw, mid, xs[1], mid, color=INK, w=2.5)
    text(s, xs[0] + bw - 0.1, mid - 0.42, 0.99, 0.32, "attack", size=9.5, color=GRAY,
         align="c", after=0)
    arrow(s, xs[1] + bw, mid, xs[2], mid, color=INK, w=2.5)
    text(s, xs[1] + bw - 0.1, mid - 0.42, 0.99, 0.32, "reply", size=9.5, color=GRAY,
         align="c", after=0)
    # verdict split
    ref_r = xs[2] + bw          # 7.35
    vx, ow, oh = 8.05, 4.82, 1.05
    win_y, loss_y = 1.95, 3.65
    node(s, vx, win_y, ow, oh, "Red WINS the round",
         sub="→ Test Manager regression, filed automatically", fill=WHITE,
         tcolor=CYAN, scolor=INKSOFT, line=CYAN, line_w=1.75, tsize=15, ssize=11)
    node(s, vx, loss_y, ow, oh, "Blue HOLDS the line",
         sub="→ Coach updates priors, invents a harder attack", fill=WHITE,
         tcolor=TEAL, scolor=INKSOFT, line=TEAL, line_w=1.75, tsize=15, ssize=11)
    arrow(s, ref_r, mid, ref_r + 0.4, mid, color=INK, w=2.5, tail=None)
    arrow(s, ref_r + 0.4, mid, vx, win_y + oh / 2, color=CYAN, w=2.0, head="triangle",
          tail=None)
    arrow(s, ref_r + 0.4, mid, vx, loss_y + oh / 2, color=TEAL, w=2.0, head="triangle",
          tail=None)
    # the loop-back
    text(s, 0.42, 5.5, 12.4, 0.9,
         [[{"t": "The corpus is the product.  ", "size": 17, "bold": True,
            "color": INK},
           {"t": "Because the Coach is rewarded for landing attacks, the regression "
            "suite grows itself. Every fight makes the next agent harder to break.",
            "size": 17, "color": INKSOFT}]], spacing=1.12, after=0)


def s08_redblue(d):
    s = content_slide(d, "Red team and blue team, in one arena", "03", "What it is")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "Two disciplines borrowed from security. Gauntlet automates both and runs "
         "them against each other, continuously.", size=16, color=INKSOFT, after=0)
    # ── RED source (attacker) on the left ──
    rect(s, 0.72, 2.95, 1.25, 1.95, fill=ORANGE, shape=MSO_SHAPE.LIGHTNING_BOLT)
    chip(s, 0.5, 5.05, 1.7, "RED TEAM", fill=ORANGE, size=11)
    text(s, 0.5, 5.5, 1.7, 0.4, "breaks it", size=12, italic=True, color=RUST,
         align="c", after=0)
    # ── attack vectors: labelled arrows flying at the shield ──
    lanes = [(3.02, "Prompt injection, direct and indirect"),
             (3.92, "Impersonation and social pressure"),
             (4.82, "Data exfiltration and policy bypass")]
    for ly, lab in lanes:
        rect(s, 2.35, ly, 7.75, 0.66, fill=ORANGE, shape=MSO_SHAPE.RIGHT_ARROW)
        text(s, 2.6, ly, 6.7, 0.66, lab, size=12.5, bold=True, color=WHITE, anchor="m",
             after=0)
    # ── BLUE shield (defender) on the right ──
    rect(s, 10.32, 4.55, 2.5, 1.05, fill=TEAL, shape=MSO_SHAPE.ISOSCELES_TRIANGLE,
         rot=180)
    rect(s, 10.32, 2.78, 2.5, 1.95, fill=TEAL, radius=0.08)
    text(s, 10.32, 2.95, 2.5, 0.55, "✓", size=30, bold=True, color=WHITE, align="c",
         after=0)
    text(s, 10.32, 3.62, 2.5, 0.5, "BLUE TEAM", size=15, bold=True, color=WHITE,
         align="c", after=0)
    text(s, 10.32, 4.08, 2.5, 0.5, "holds the line", size=11.5, italic=True,
         color=LIGHTCYAN, align="c", after=0)
    # ── the takeaway ──
    card(s, 0.42, 6.02, 12.5, 0.82, fill=MIST, radius=0.05)
    text(s, 0.75, 6.02, 12.0, 0.82,
         [[{"t": "Blue defends with input distrust, identity checks and escalation; "
            "an independent Referee agent scores every exchange. ", "size": 13,
            "color": INKSOFT},
           {"t": "Gauntlet plays both sides across 8 personas and 4 postures.",
            "size": 13, "bold": True, "color": INK}]], anchor="m", spacing=1.12,
         after=0)


def s09_thesis(d):
    s = content_slide(d, "It takes an adversarial AI to harden an AI", "03",
                      "The thesis")
    text(s, 0.42, 1.55, 7.4, 3.0,
         [[{"t": "You cannot hand write every attack.\n", "size": 20, "bold": True,
            "color": INK}],
          [{"t": "The space of adversarial prompts is effectively infinite and it "
            "shifts with every model update. A fixed checklist is out of date the day "
            "you write it.", "size": 15.5, "color": INKSOFT}],
          [{"t": "\n", "size": 8}],
          [{"t": "So Gauntlet does what AlphaGo did: ", "size": 15.5, "color": INKSOFT},
           {"t": "self play.", "size": 15.5, "bold": True, "color": ORANGE}],
          [{"t": "An AI adversary that learns, mutates, and escalates is the only "
            "thing that scales to an AI defender. Safety becomes a moving number you "
            "measure, not a box you tick.", "size": 15.5, "color": INKSOFT}]],
         spacing=1.16, after=6)
    # right: escalation ladder
    lx, lw = 8.35, 4.55
    card(s, lx, 1.75, lw, 4.7, fill=INK, radius=0.04)
    text(s, lx + 0.35, 2.0, lw - 0.7, 0.4, "SELF PLAY ESCALATION", size=12, bold=True,
         color=ORANGE, after=0)
    rungs = [
        ("Known corpus", "replay proven attacks", CORAL),
        ("Mutation", "recombine what nearly worked", ORANGE),
        ("Invention", "author a brand new persona", LIGHTCYAN),
        ("Harder agent", "every win is a new test", CYAN),
    ]
    for i, (t_, sub, c) in enumerate(rungs):
        yy = 5.55 - i * 0.92
        rect(s, lx + 0.35, yy, 0.14, 0.72, fill=c)
        text(s, lx + 0.62, yy - 0.02, lw - 1.0, 0.4, t_, size=15, bold=True,
             color=WHITE, after=0)
        text(s, lx + 0.62, yy + 0.34, lw - 1.0, 0.35, sub, size=11, color=GRAY,
             after=0)
        if i < 3:
            text(s, lx + 0.4, yy - 0.32, 0.5, 0.3, "↑", size=15, bold=True,
                 color=c, after=0)


def s10_section_how(d):
    s = d.add("Section Header")
    rect(s, -0.1, -0.1, d.W + 0.2, d.H + 0.2, fill=INK)
    rect(s, 0.7, 2.35, 1.9, 0.11, fill=ORANGE)
    text(s, 0.68, 2.7, 11.0, 0.55, "STEP BY STEP", size=15, bold=True, color=ORANGE,
         after=0)
    text(s, 0.66, 3.3, 12.0, 1.7, "How a fight works", size=60, bold=True, color=WHITE,
         after=0)
    text(s, 0.7, 5.0, 11.2, 0.8,
         "The cast, one round end to end, and the loop that never runs out of attacks.",
         size=18, color=RGBColor(0xC9, 0xCE, 0xD2), after=0)


def s11_cast(d):
    s = content_slide(d, "The cast: four agents, four jobs", "04", "How it works")
    cast = [
        ("Red Coach", "Attacker", "Picks the highest reward attack or invents a new "
         "one. Python + LangGraph.", ORANGE),
        ("Blue target", "Defender", "The agent under test. Agent Builder agent, or any "
         "external target.", TEAL),
        ("Referee", "Judge", "Scores each transcript against the rubric. Independent "
         "Agent Builder agent.", PURPLE),
        ("Fix Recommender", "Medic", "On a loss, drafts root cause + prompt patch + "
         "regression test.", CYAN),
    ]
    x, w, gap = 0.42, 3.05, 0.16
    for i, (name, role, desc, c) in enumerate(cast):
        lx = x + i * (w + gap)
        card(s, lx, 2.05, w, 4.35, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.04)
        rect(s, lx, 2.05, w, 1.0, fill=c, radius=0.04)
        rect(s, lx, 2.55, w, 0.5, fill=c)  # square off bottom of header
        text(s, lx + 0.25, 2.05, w - 0.5, 1.0,
             [[{"t": role.upper(), "size": 11, "bold": True,
                "color": RGBColor(0xFF, 0xFF, 0xFF)}]], anchor="m", after=0)
        text(s, lx + 0.25, 3.25, w - 0.5, 0.8, name, size=17, bold=True, color=INK,
             after=0)
        text(s, lx + 0.25, 4.15, w - 0.5, 1.9, desc, size=12.5, color=INKSOFT,
             spacing=1.12, after=0)
        text(s, lx + 0.25, 6.0, w - 0.5, 0.35,
             ["Agent Builder", "Agent Builder", "Agent Builder",
              "Coded agent"][i] if i != 0 else "Coded agent", size=10, bold=True,
             color=c, after=0)
    footer(s, "Deliberately polyglot: coded agents and low code Agent Builder agents "
              "fight in the same ring.")


def s12_round(d):
    s = content_slide(d, "One round, end to end", "04", "How it works")
    y = 2.5
    bw, bh, gap = 2.55, 1.15, 0.55
    stages = [("Matched", "persona vs\nposture", TEAL),
              ("Negotiating", "multi turn\nRed ↔ Blue", ORANGE),
              ("Verdict", "Referee\nscores it", PURPLE),
              ("Archived", "outcome\nrouted", INK)]
    xs = []
    for i, (t_, sub, c) in enumerate(stages):
        lx = 0.42 + i * (bw + gap)
        xs.append(lx)
        node(s, lx, y, bw, bh, t_, sub=sub, fill=WHITE, tcolor=c, scolor=INKSOFT,
             line=c, line_w=1.75, tsize=17, ssize=11)
        if i < 3:
            arrow(s, lx + bw + 0.04, y + bh / 2, lx + bw + gap - 0.04, y + bh / 2,
                  color=INK, w=2.5)
    # negotiating loop marker
    text(s, xs[1], y + bh + 0.08, bw, 0.3, "↺ up to N turns", size=10.5,
         italic=True, color=ORANGE, align="c", after=0)
    # outcomes below Archived
    ax = xs[3]
    win = node(s, ax - 0.6, 4.6, 3.75, 0.85,
               "Red won → Test Manager", sub="new regression case", fill=WHITE,
               tcolor=CYAN, scolor=INKSOFT, line=CYAN, line_w=1.5, tsize=13, ssize=10.5)
    node(s, ax - 0.6, 5.65, 3.75, 0.85, "Blue lost → Action Center",
         sub="fix task, patch attached", fill=WHITE, tcolor=RUST, scolor=INKSOFT,
         line=RUST, line_w=1.5, tsize=13, ssize=10.5)
    arrow(s, ax + bw / 2, y + bh, ax + bw / 2, 4.6, color=INK, w=2.0, tail=None,
          head="triangle")
    # steps list on the left-bottom
    text(s, 0.42, 4.5, 7.0, 0.4, "What RoundOrchestrator (Maestro Flow) does",
         size=13, bold=True, color=INK, after=0)
    for i, st in enumerate([
        "Pull the next Red persona from the corpus (or ask the Coach to invent one).",
        "Run N exchanges until Red gives up or Blue capitulates.",
        "Send the transcript + tool ledger to the Referee.",
        "Post the verdict to the Case; on a Blue loss, trigger the Fix Recommender.",
    ]):
        yy = 5.0 + i * 0.44
        text(s, 0.42, yy, 7.2, 0.42,
             [[{"t": f"{i + 1}.  ", "bold": True, "color": ORANGE, "size": 12},
               {"t": st, "size": 12, "color": INKSOFT}]], after=0)


def s13_coach(d):
    s = content_slide(d, "The Coach that never runs out of attacks", "04",
                      "How it works")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "A self play loop over the corpus. It spends its attempts where they are most "
         "likely to break the target.", size=16, color=INKSOFT, after=0)
    steps = [
        ("1", "Pick", "Choose the persona with the highest risk weighted expected "
         "reward against this target (Thompson sampling over the corpus).", ORANGE),
        ("2", "Invent", "If nothing has positive expected reward, ask the LLM to author "
         "a new persona, conditioned on the last K losing transcripts.", PURPLE),
        ("3", "Fight", "Run the round through RoundOrchestrator and let the Referee "
         "score it.", TEAL),
        ("4", "Learn", "Win → persist as a Test Manager regression. Loss → "
         "update reward priors. Then loop.", CYAN),
    ]
    x, w, gap = 0.42, 3.05, 0.16
    for i, (n, h, b, c) in enumerate(steps):
        lx = x + i * (w + gap)
        card(s, lx, 2.35, w, 3.5, fill=MIST, radius=0.04)
        rect(s, lx + 0.28, 2.62, 0.62, 0.62, fill=c, radius=0.5)
        text(s, lx + 0.28, 2.63, 0.62, 0.62, n, size=24, bold=True, color=WHITE,
             align="c", anchor="m", after=0)
        text(s, lx + 0.28, 3.45, w - 0.5, 0.5, h, size=17, bold=True, color=INK,
             after=0)
        text(s, lx + 0.28, 4.05, w - 0.55, 1.7, b, size=12.5, color=INKSOFT,
             spacing=1.12, after=0)
        if i < 3:
            text(s, lx + w - 0.02, 3.7, 0.2, 0.4, "›", size=22, bold=True,
                 color=GRAY, align="c", after=0)
    # loop-back ribbon
    rect(s, 0.42, 6.05, 12.5, 0.5, fill=INK, radius=0.06)
    text(s, 0.42, 6.05, 12.5, 0.5,
         [[{"t": "↺  ", "size": 15, "bold": True, "color": ORANGE},
           {"t": "No human writes the next test. The corpus is the byproduct of the "
            "Coach trying to win.", "size": 13.5, "color": WHITE, "bold": True}]],
         align="c", anchor="m", after=0)


def s14_injection(d):
    s = content_slide(d, "A real attack: injection through the data, not the user",
                      "04", "How it works")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "The scariest attacks never touch the chat box. They ride in on data your "
         "agent is trusted to read.", size=16, color=INKSOFT, after=0)
    # flow: attacker -> CRM -> tool result -> Blue
    y = 2.5
    node(s, 0.42, y, 2.3, 1.1, "Attacker", sub="plants a note", fill=RUST, tcolor=WHITE,
         scolor=CORAL, tsize=15, ssize=10.5)
    node(s, 3.35, y, 2.6, 1.1, "CRM record", sub="“KYC subsystem” memo",
         fill=WHITE, tcolor=INK, scolor=INKSOFT, line=GRAY, line_w=1.5, tsize=15,
         ssize=10.5)
    node(s, 6.55, y, 3.0, 1.1, "fetch_customer_\ncorrespondence", sub="tool result",
         fill=MIST, tcolor=INK, scolor=INKSOFT, line=ORANGE, line_w=1.75, tsize=12.5,
         ssize=10.5)
    node(s, 10.15, y, 2.75, 1.1, "Blue agent", sub="reads it as trusted",
         fill=TEAL, tcolor=WHITE, scolor=LIGHTCYAN, tsize=15, ssize=10.5)
    for x1, x2 in [(2.72, 3.35), (5.95, 6.55), (9.55, 10.15)]:
        arrow(s, x1, y + 0.55, x2, y + 0.55, color=INK, w=2.5)
    text(s, 6.55, y + 1.16, 3.0, 0.3, "not the user channel", size=10, italic=True,
         color=ORANGE, align="c", after=0)
    # outcome banner
    card(s, 0.42, 4.35, 12.5, 1.05, fill=INK, radius=0.04)
    rect(s, 0.42, 4.35, 0.14, 1.05, fill=ORANGE)
    text(s, 0.8, 4.35, 12.0, 1.05,
         [[{"t": "Headline breach:  ", "size": 16, "bold": True, "color": ORANGE},
           {"t": "a planted “KYC subsystem” note talked a naive agent into "
            "disclosing an ", "size": 15, "color": WHITE},
           {"t": "$84,320 balance", "size": 16, "bold": True, "color": WHITE},
           {"t": " with ", "size": 15, "color": WHITE},
           {"t": "zero 2FA.", "size": 16, "bold": True, "color": CORAL}]],
         anchor="m", spacing=1.15, after=0)
    text(s, 0.42, 5.75, 12.4, 0.7,
         "This is OWASP LLM01 (Prompt Injection) arriving as OWASP LLM02 (Sensitive "
         "Information Disclosure). Gauntlet seeds these payloads on purpose, through "
         "the tool layer, exactly where real ones hide.", size=13.5, color=INKSOFT,
         spacing=1.14, after=0)


def s15_flywheel(d):
    s = content_slide(d, "Every outcome does work: the flywheel", "04", "How it works")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "A fight is never just a score. Each verdict writes back into UiPath and makes "
         "the system stronger.", size=16, color=INKSOFT, after=0)
    # center fight node
    cx, cy = 5.9, 3.7
    rect(s, cx, cy, 2.5, 1.15, fill=INK, radius=0.06)
    text(s, cx, cy, 2.5, 1.15, [[{"t": "FIGHT", "size": 18, "bold": True,
         "color": WHITE}], [{"t": "Red vs Blue", "size": 11, "color": GRAY}]],
         align="c", anchor="m", after=2)
    # win branch (up-right)
    node(s, 9.4, 2.15, 3.5, 1.0, "WIN → Test Manager",
         sub="permanent regression test, auto filed", fill=WHITE, tcolor=CYAN,
         scolor=INKSOFT, line=CYAN, line_w=1.75, tsize=14, ssize=11)
    arrow(s, cx + 2.5, cy + 0.2, 9.4, 2.65, color=CYAN, w=2.5)
    # loss branch (down-right)
    node(s, 9.4, 5.4, 3.5, 1.0, "LOSS → Action Center",
         sub="fix task with patch already drafted", fill=WHITE, tcolor=RUST,
         scolor=INKSOFT, line=RUST, line_w=1.75, tsize=14, ssize=11)
    arrow(s, cx + 2.5, cy + 0.95, 9.4, 5.9, color=RUST, w=2.5)
    # harder agent -> back to fight (left loop)
    node(s, 0.42, 3.75, 3.5, 1.05, "Hardened agent", sub="+ bigger corpus", fill=ORANGE,
         tcolor=WHITE, scolor=RGBColor(0xFF, 0xD9, 0xCC), tsize=15, ssize=11)
    arrow(s, 3.92, 4.05, cx, cy + 0.35, color=ORANGE, w=2.5)
    # feed lines back
    text(s, 0.42, 6.35, 12.5, 0.6,
         [[{"t": "The result: ", "size": 15, "bold": True, "color": INK},
           {"t": "a regression suite that grows itself and an auditor ready coverage "
            "matrix, a byproduct of running fights.", "size": 15,
            "color": INKSOFT}]], after=0)


def s16_section_uipath(d):
    s = d.add("Section Header")
    rect(s, -0.1, -0.1, d.W + 0.2, d.H + 0.2, fill=INK)
    rect(s, 0.7, 2.35, 1.9, 0.11, fill=ORANGE)
    text(s, 0.68, 2.7, 11.0, 0.55, "NOT A PROTOTYPE · DEPLOYED", size=15,
         bold=True, color=ORANGE, after=0)
    text(s, 0.66, 3.3, 12.0, 1.7, "Built on UiPath", size=64, bold=True, color=WHITE,
         after=0)
    text(s, 0.7, 5.0, 11.2, 0.8,
         "Maestro, Agent Builder, Test Manager, Action Center, Coded Apps, live "
         "in the Labs tenant, not slideware.", size=18,
         color=RGBColor(0xC9, 0xCE, 0xD2), after=0)


def s17_architecture(d):
    s = content_slide(d, "Architecture: an agentic test factory on UiPath", "05",
                      "Architecture")
    # ── Plane 1: Console (top) ──
    rect(s, 0.42, 1.5, 12.5, 0.86, fill=INK, radius=0.04)
    text(s, 0.62, 1.5, 3.4, 0.86, [[{"t": "Coded App", "size": 13, "bold": True,
         "color": WHITE}], [{"t": "gauntletapp · React", "size": 10,
         "color": GRAY}]], anchor="m", after=0)
    for i, t_ in enumerate(["Dashboard", "Coach Lab", "Fix Lab", "Analytics", "Audit",
                            "Logs"]):
        lx = 4.05 + i * 1.46
        rect(s, lx, 1.66, 1.34, 0.54, fill=RGBColor(0x2B, 0x36, 0x3B), radius=0.1)
        text(s, lx, 1.66, 1.34, 0.54, t_, size=10.5, color=WHITE, align="c",
             anchor="m", after=0)
    text(s, 0.42, 2.42, 12.5, 0.3,
         "browser side calls via @uipath/uipath-typescript SDK  ↓", size=9.5,
         italic=True, color=GRAY, align="c", after=0)
    # ── Plane 2: Automation Cloud / Maestro (middle) ──
    rect(s, 0.42, 2.85, 8.5, 3.0, fill=MIST, line=CLOUD, line_w=1.25, radius=0.03)
    text(s, 0.62, 2.95, 8.2, 0.3,
         "UiPath Automation Cloud · Orchestrator folder Shared/Gauntlet", size=10.5,
         bold=True, color=TEAL, after=0)
    # FightArena case wrapping RoundOrchestrator
    rect(s, 0.7, 3.35, 7.95, 2.35, fill=WHITE, line=TEAL, line_w=1.5, radius=0.03)
    text(s, 0.85, 3.42, 5.0, 0.3, "FightArena · Maestro Case", size=11, bold=True,
         color=TEAL, after=0)
    rect(s, 0.9, 3.8, 7.55, 1.78, fill=RGBColor(0xE9, 0xF1, 0xF4), radius=0.03)
    text(s, 1.02, 3.85, 6.0, 0.3, "RoundOrchestrator · Maestro Flow (per round)",
         size=10.5, bold=True, color=TEAL, after=0)
    ay = 4.25
    node(s, 1.05, ay, 1.62, 1.05, "Red Coach", sub="Python", fill=ORANGE, tcolor=WHITE,
         scolor=RGBColor(0xFF, 0xD9, 0xCC), tsize=12, ssize=9.5)
    node(s, 2.87, ay, 1.62, 1.05, "MetroBankCSR", sub="Agent Builder", fill=TEAL,
         tcolor=WHITE, scolor=LIGHTCYAN, tsize=11, ssize=9)
    node(s, 4.69, ay, 1.5, 1.05, "RefereeAgent", sub="Agent Builder", fill=PURPLE,
         tcolor=WHITE, scolor=RGBColor(0xE7, 0xC4, 0xE6), tsize=10.5, ssize=9)
    node(s, 6.39, ay, 1.85, 1.05, "Fix Recommender", sub="Python", fill=CYAN,
         tcolor=WHITE, scolor=RGBColor(0xD6, 0xF2, 0xF5), tsize=11, ssize=9.5)
    arrow(s, 2.67, ay + 0.5, 2.87, ay + 0.5, color=INK, w=1.75)
    arrow(s, 4.49, ay + 0.5, 4.69, ay + 0.5, color=INK, w=1.75)
    # ── Plane 3: Outcomes (right) ──
    rect(s, 9.1, 2.85, 3.82, 3.0, fill=INK, radius=0.03)
    text(s, 9.3, 2.95, 3.5, 0.3, "WRITES BACK TO YOUR TENANT", size=10, bold=True,
         color=ORANGE, after=0)
    node(s, 9.32, 3.4, 3.4, 1.0, "Test Manager", sub="regressions grow themselves",
         fill=WHITE, tcolor=CYAN, scolor=GRAY, line=CYAN, line_w=1.5, tsize=13,
         ssize=9.5)
    node(s, 9.32, 4.55, 3.4, 1.0, "Action Center", sub="human approved fixes",
         fill=WHITE, tcolor=CORAL, scolor=GRAY, line=CORAL, line_w=1.5, tsize=13,
         ssize=9.5)
    arrow(s, 8.92, 4.3, 9.32, 3.9, color=CYAN, w=2.0)
    arrow(s, 8.92, 4.7, 9.32, 5.05, color=CORAL, w=2.0)
    # ── Plane 4: corpus ──
    rect(s, 0.42, 6.05, 12.5, 0.62, fill=RGBColor(0xFF, 0xEC, 0xE6), line=ORANGE,
         line_w=1.25, radius=0.06)
    text(s, 0.42, 6.05, 12.5, 0.62,
         [[{"t": "runs/  ·  the shared corpus  ", "size": 12.5, "bold": True,
            "color": RUST},
           {"t": "read by the Console, Test Manager import, and the offline dashboard. "
            "The single source of truth.", "size": 12, "color": INKSOFT}]],
         align="c", anchor="m", after=0)


def s18_components(d):
    s = content_slide(d, "Framework neutral by design", "05", "Architecture")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "The arena does not care how the target is built. Coded, low code, and "
         "third party agents all fight under the same rules.", size=16, color=INKSOFT,
         after=0)
    groups = [
        ("LOW CODE", TEAL, [
            "MetroBankCSR · Agent Builder (Blue target)",
            "RefereeAgent · Agent Builder (judge)",
            "FightArena · Maestro Case",
            "RoundOrchestrator · Maestro Flow"]),
        ("CODED", ORANGE, [
            "Red Coach · Python + LangGraph",
            "Fix Recommender · Python + LangGraph",
            "Referee mirror · local judge logic",
            "gauntletapp · React + TypeScript SDK"]),
        ("EXTERNAL", PURPLE, [
            "LangGraph target · third party agent",
            "No UiPath tool access, still testable",
            "Proves the arena is not UiPath only",
            "Same referee, same corpus, same tags"]),
    ]
    x, w, gap = 0.42, 4.05, 0.17
    for i, (name, c, items) in enumerate(groups):
        lx = x + i * (w + gap)
        card(s, lx, 2.35, w, 4.0, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.04)
        rect(s, lx, 2.35, w, 0.6, fill=c, radius=0.04)
        rect(s, lx, 2.7, w, 0.25, fill=c)
        text(s, lx, 2.35, w, 0.6, name, size=13, bold=True, color=WHITE, align="c",
             anchor="m", after=0)
        for j, it in enumerate(items):
            yy = 3.2 + j * 0.72
            rect(s, lx + 0.28, yy + 0.06, 0.14, 0.14, fill=c, radius=0.5)
            text(s, lx + 0.54, yy, w - 0.8, 0.66, it, size=12, color=INK, spacing=1.05,
                 after=0)
    footer(s, "One target contract. Eight UiPath products. Zero lock in on how your "
              "agent is authored.")


def _img_aspect(path):
    if _PILImage is not None:
        try:
            with _PILImage.open(str(path)) as im:
                return im.size[0] / im.size[1]
        except Exception:
            pass
    return 1.6


def place_fit(slide, path, bl, bt, bw, bh):
    """Draw the image fit inside (bl,bt,bw,bh), preserving aspect and centring."""
    if not path.exists():
        return
    a = _img_aspect(path)
    if bw / bh > a:                       # box wider than image -> height limited
        dh, dw = bh, bh * a
    else:                                 # width limited
        dw, dh = bw, bw / a
    slide.shapes.add_picture(str(path), Inches(bl + (bw - dw) / 2),
                             Inches(bt + (bh - dh) / 2), width=Inches(dw),
                             height=Inches(dh))


def browser_frame(slide, l, t, w, h, url):
    """A clean app-window chrome; returns the inner box for the screenshot."""
    card(slide, l, t, w, h, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.02)
    rect(slide, l, t, w, 0.34, fill=CLOUD, radius=0.02)
    rect(slide, l, t + 0.17, w, 0.17, fill=CLOUD)  # square the bar's lower corners
    for i in range(3):
        rect(slide, l + 0.2 + i * 0.22, t + 0.11, 0.12, 0.12,
             fill=RGBColor(0xB8, 0xBE, 0xC3), radius=0.5)
    rect(slide, l + 1.05, t + 0.06, w - 1.5, 0.22, fill=WHITE, radius=0.5)
    text(slide, l + 1.25, t + 0.05, w - 1.7, 0.24, url, size=8.5, color=GRAY,
         anchor="m", after=0)
    return (l + 0.16, t + 0.44, w - 0.32, h - 0.58)


def app_slide(d, title, shot, step, label, bullets):
    """One in-app walkthrough slide: big framed screenshot + step + captions."""
    s = content_slide(d, title, "05", "In the app")
    inner = browser_frame(s, 0.42, 1.74, 8.72, 5.15,
                          "hackathon26_038.staging.uipath.host/gauntletapp")
    place_fit(s, SHOTS / shot, *inner)
    rx = 9.4
    rect(s, rx, 1.95, 0.66, 0.66, fill=ORANGE, radius=0.14)
    text(s, rx, 1.95, 0.66, 0.66, step, size=27, bold=True, color=WHITE, align="c",
         anchor="m", after=0)
    text(s, rx + 0.86, 1.98, 3.0, 0.62, label, size=20, bold=True, color=INK,
         anchor="m", after=0)
    rect(s, rx, 2.85, 0.55, 0.07, fill=ORANGE)
    for i, b in enumerate(bullets):
        yy = 3.2 + i * 0.98
        rect(s, rx, yy + 0.04, 0.16, 0.16, fill=ORANGE, radius=0.5)
        text(s, rx + 0.3, yy - 0.03, 3.35, 0.92, b, size=13, color=INKSOFT,
             spacing=1.16, after=0)
    footer(s, "Live tenant data via @uipath/uipath-typescript; falls back to the "
              "bundled 75 fight corpus offline.")


def s20_section_audit(d):
    s = d.add("Section Header")
    rect(s, -0.1, -0.1, d.W + 0.2, d.H + 0.2, fill=INK)
    rect(s, 0.7, 2.35, 1.9, 0.11, fill=ORANGE)
    text(s, 0.68, 2.7, 11.0, 0.55, "COVERAGE YOU CAN AUDIT", size=15, bold=True,
         color=ORANGE, after=0)
    text(s, 0.66, 3.3, 12.4, 1.7, "Mapped to the frameworks", size=52, bold=True,
         color=WHITE, after=0)
    text(s, 0.7, 5.0, 11.6, 0.8,
         "Every fight is double tagged, so “are we safe?” becomes a matrix, "
         "not an opinion. OWASP · MITRE ATLAS · ISO 42001.", size=18,
         color=RGBColor(0xC9, 0xCE, 0xD2), after=0)


def s21_owasp(d):
    s = content_slide(d, "OWASP LLM Top 10: where Blue is tested", "06", "Frameworks")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "The industry standard list of the ten worst LLM application risks. Gauntlet "
         "attacks the ones an agent actually faces.", size=15.5, color=INKSOFT,
         after=0)
    items = [
        ("LLM01", "Prompt Injection", 3), ("LLM02", "Sensitive Info Disclosure", 3),
        ("LLM06", "Excessive Agency", 3), ("LLM07", "System Prompt Leakage", 2),
        ("LLM05", "Improper Output Handling", 2), ("LLM09", "Misinformation", 1),
        ("LLM04", "Data & Model Poisoning", 1), ("LLM08", "Vector Weaknesses", 0),
        ("LLM03", "Supply Chain", 0), ("LLM10", "Unbounded Consumption", 1),
    ]
    heat = {3: ORANGE, 2: CORAL, 1: RGBColor(0xF6, 0xD0, 0xB8), 0: CLOUD}
    tcol = {3: WHITE, 2: WHITE, 1: RUST, 0: GRAY}
    x, w, gap = 0.42, 2.44, 0.13
    for i, (code, name, lvl) in enumerate(items):
        col = i % 5
        row = i // 5
        lx = x + col * (w + gap)
        ty = 2.35 + row * 1.9
        card(s, lx, ty, w, 1.72, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.05)
        rect(s, lx, ty, w, 0.5, fill=heat[lvl], radius=0.05)
        rect(s, lx, ty + 0.28, w, 0.22, fill=heat[lvl])
        text(s, lx, ty, w, 0.5, code, size=13, bold=True, color=tcol[lvl], align="c",
             anchor="m", after=0)
        text(s, lx + 0.15, ty + 0.6, w - 0.3, 0.7, name, size=12, bold=True, color=INK,
             align="c", after=0)
        label = {3: "core", 2: "covered", 1: "touched", 0: "n/a"}[lvl]
        text(s, lx + 0.15, ty + 1.3, w - 0.3, 0.3, label.upper(), size=9.5, bold=True,
             color=heat[lvl] if lvl else GRAY, align="c", after=0)
    # legend
    lx = 0.42
    for lbl, lvl in [("Core focus", 3), ("Covered", 2), ("Touched", 1),
                     ("Out of scope", 0)]:
        rect(s, lx, 6.35, 0.24, 0.24, fill=heat[lvl], radius=0.3)
        text(s, lx + 0.32, 6.32, 1.9, 0.3, lbl, size=10.5, color=INKSOFT, after=0)
        lx += 2.5


def s22_mitre(d):
    s = content_slide(d, "MITRE ATLAS: the adversary's playbook", "06", "Frameworks")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "ATLAS maps real attacks on AI systems to tactics and techniques, the "
         "ATT&CK of machine learning. Gauntlet tags each fight to it.", size=15,
         color=INKSOFT, after=0)
    tactics = [
        ("Reconnaissance", "probe the agent's role & limits", GRAY),
        ("Initial Access", "AML.T0051 · LLM Prompt Injection", ORANGE),
        ("ML Attack Staging", "AML.T0054 · LLM Jailbreak", RUST),
        ("Exfiltration", "AML.T0057 · LLM Data Leakage", PURPLE),
        ("Impact", "unauthorised action / disclosure", INK),
    ]
    x, w, gap = 0.42, 2.4, 0.1
    y = 2.6
    for i, (name, tech, c) in enumerate(tactics):
        lx = x + i * (w + gap)
        card(s, lx, y, w, 2.5, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.04)
        rect(s, lx, y, w, 0.12, fill=c)
        text(s, lx + 0.18, y + 0.3, w - 0.36, 0.3, f"TACTIC {i + 1}", size=9.5,
             bold=True, color=c, after=0)
        text(s, lx + 0.18, y + 0.66, w - 0.36, 0.9, name, size=14.5, bold=True,
             color=INK, spacing=1.02, after=0)
        text(s, lx + 0.18, y + 1.62, w - 0.36, 0.8, tech, size=11, color=INKSOFT,
             spacing=1.1, after=0)
        if i < 4:
            arrow(s, lx + w + 0.005, y + 1.25, lx + w + gap - 0.005, y + 1.25,
                  color=GRAY, w=1.75)
    card(s, 0.42, 5.55, 12.5, 0.95, fill=MIST, radius=0.04)
    text(s, 0.7, 5.55, 12.0, 0.95,
         [[{"t": "Why two frameworks?  ", "size": 14, "bold": True, "color": INK},
           {"t": "OWASP names the vulnerability class; ATLAS names the adversary "
            "behaviour. Double tagging gives a security team the language it already "
            "uses on its other systems.", "size": 13.5, "color": INKSOFT}]],
         anchor="m", spacing=1.14, after=0)


def s23_iso(d):
    s = content_slide(d, "The evidence engine for AI governance", "06", "Frameworks")
    text(s, 0.42, 1.5, 12.4, 0.6,
         "Standards and regulations increasingly demand adversarial testing and "
         "records of it. Gauntlet produces that evidence as a byproduct of running.",
         size=15.5, color=INKSOFT, after=0)
    cards = [
        ("ISO/IEC 42001", "AI Management System",
         "Requires risk assessment, controls, and continual monitoring of AI systems. "
         "Gauntlet is the recurring control + the audit trail.", TEAL),
        ("NIST AI RMF", "Measure & Manage",
         "The Measure function calls for red teaming and resilience testing. Gauntlet "
         "operationalises it and logs every run.", CYAN),
        ("EU AI Act", "Art. 15 · robustness",
         "High risk & GPAI systems must be robust against adversarial manipulation. "
         "Gauntlet is continuous proof of that testing.", PURPLE),
    ]
    x, w, gap = 0.42, 4.05, 0.17
    for i, (name, tag, body, c) in enumerate(cards):
        lx = x + i * (w + gap)
        card(s, lx, 2.4, w, 3.35, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.04)
        rect(s, lx, 2.4, 0.12, 3.35, fill=c)
        text(s, lx + 0.35, 2.65, w - 0.6, 0.5, name, size=18, bold=True, color=INK,
             after=0)
        text(s, lx + 0.35, 3.2, w - 0.6, 0.35, tag.upper(), size=10.5, bold=True,
             color=c, after=0)
        text(s, lx + 0.35, 3.75, w - 0.65, 1.9, body, size=12.5, color=INKSOFT,
             spacing=1.16, after=0)
    card(s, 0.42, 6.0, 12.5, 0.62, fill=INK, radius=0.06)
    text(s, 0.42, 6.0, 12.5, 0.62,
         [[{"t": "Gauntlet does not certify you. ", "size": 13.5, "bold": True,
            "color": ORANGE},
           {"t": "It gives you the artifact an auditor asks for: proof you attacked "
            "your own agent, and what happened.", "size": 13.5, "color": WHITE}]],
         align="c", anchor="m", after=0)


def s24_proof(d):
    s = content_slide(d, "The proof: measured, not asserted", "07", "Results")
    kpi(s, 0.42, 1.7, 3.02, "75", "adversarial fights run across personas, scenarios "
        "and postures", accent=INK, h=1.7)
    kpi(s, 3.62, 1.7, 3.02, "27/27", "attacks blocked under the Standard (strict) "
        "posture", accent=TEAL, h=1.7)
    kpi(s, 6.82, 1.7, 3.02, "90.6", "average defense score out of 100 at strict "
        "posture", accent=CYAN, h=1.7)
    kpi(s, 10.02, 1.7, 2.9, "22%", "worst case attacker win rate once the policy is "
        "weakened", accent=ORANGE, h=1.7)
    # headline breach band
    card(s, 0.42, 3.75, 8.35, 2.65, fill=INK, radius=0.04)
    rect(s, 0.42, 3.75, 0.14, 2.65, fill=ORANGE)
    text(s, 0.85, 3.95, 7.6, 0.4, "THE HEADLINE BREACH", size=12, bold=True,
         color=ORANGE, after=0)
    text(s, 0.85, 4.4, 7.7, 1.9,
         [[{"t": "$84,320", "size": 44, "bold": True, "color": WHITE}],
          [{"t": "disclosed by a naive agent with zero 2FA, after a planted "
            "“KYC subsystem” note convinced it the request was internal.",
            "size": 14.5, "color": RGBColor(0xC9, 0xCE, 0xD2)}]], spacing=1.12,
         after=6)
    # right column: what it produced
    rx = 9.0
    card(s, rx, 3.75, 3.92, 2.65, fill=MIST, radius=0.04)
    text(s, rx + 0.3, 3.95, 3.4, 0.4, "WHAT IT PRODUCED", size=12, bold=True,
         color=TEAL, after=0)
    for i, (n, lab) in enumerate([("9", "critical findings auto diagnosed, each with a "
                                   "patch proposed"),
                                  ("∞", "a regression suite that grows itself"),
                                  ("1", "OWASP + MITRE matrix an auditor can read")]):
        yy = 4.4 + i * 0.66
        text(s, rx + 0.3, yy, 0.6, 0.6, n, size=22, bold=True, color=ORANGE, after=0)
        text(s, rx + 0.95, yy + 0.04, 2.85, 0.66, lab, size=11.5, color=INKSOFT,
             spacing=1.06, after=0)


def s25_live(d):
    s = content_slide(d, "Deployed, live, and open", "07", "Close")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "Not a mockup. The full solution is deployed and activated in the UiPath Labs "
         "tenant provisioned for this submission.", size=16, color=INKSOFT, after=0)
    tiles = [
        ("LABS TENANT", "staging.uipath.com/\nhackathon26_038", "Solution deployed to "
         "Shared/Gauntlet", TEAL),
        ("LIVE APP", "hackathon26_038.staging\n.uipath.host/gauntletapp", "Open it from "
         "the Apps menu", ORANGE),
        ("DEMO", "youtu.be/1q9W5SC_fxA", "Persona authored live on stage, not "
         "recorded in advance", PURPLE),
        ("CODE", "github.com/tdries/\nuipath-hackathon-gauntlet", "MIT licensed, "
         "reproducible", INK),
    ]
    x, w, gap = 0.42, 3.05, 0.16
    for i, (tag, val, sub, c) in enumerate(tiles):
        lx = x + i * (w + gap)
        card(s, lx, 2.35, w, 2.9, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.04)
        rect(s, lx, 2.35, w, 0.1, fill=c)
        text(s, lx + 0.28, 2.62, w - 0.5, 0.35, tag, size=11, bold=True, color=c,
             after=0)
        text(s, lx + 0.28, 3.05, w - 0.5, 1.1, val, size=13.5, bold=True, color=INK,
             spacing=1.05, after=0)
        text(s, lx + 0.28, 4.35, w - 0.5, 0.8, sub, size=11.5, color=INKSOFT,
             spacing=1.1, after=0)
    # vision band
    card(s, 0.42, 5.6, 12.5, 1.05, fill=INK, radius=0.04)
    rect(s, 0.42, 5.6, 0.14, 1.05, fill=ORANGE)
    text(s, 0.85, 5.6, 12.0, 1.05,
         [[{"t": "The vision:  ", "size": 16, "bold": True, "color": ORANGE},
           {"t": "continuous adversarial assurance for every agent you ship, the "
            "same way you already run unit tests and CI. An agent that keeps agents "
            "honest.", "size": 15.5, "color": WHITE}]], anchor="m", spacing=1.14,
         after=0)


def s26_close(d):
    s = d.add("Headline only")
    d.set_title(s, "")
    rect(s, -0.1, -0.1, d.W + 0.2, d.H + 0.2, fill=INK)
    rect(s, 0, 0, d.W, 0.16, fill=ORANGE)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(0.85), Inches(1.55), width=Inches(5.4))
    text(s, 0.9, 3.5, 11.6, 1.6, "Go safe, or go home.", size=62, bold=True,
         color=WHITE, after=0)
    rect(s, 0.95, 5.15, 2.7, 0.09, fill=ORANGE)
    text(s, 0.9, 5.45, 11.5, 0.6,
         "Gauntlet · Adversarial Test Cloud for UiPath agents.", size=18,
         color=RGBColor(0xC9, 0xCE, 0xD2), after=0)
    text(s, 0.9, 6.55, 11.6, 0.5,
         [[{"t": "Tim Dries · Biztory", "size": 13, "bold": True, "color": WHITE},
           {"t": "     tim.dries@biztory.be     youtu.be/1q9W5SC_fxA", "size": 13,
            "color": GRAY}]], after=0)
    text(s, 9.5, 6.95, 3.4, 0.4, "UiPath AgentHack 2026 · Track 3", size=11,
         color=GRAY, align="r", after=0)


# ── appendix ──────────────────────────────────────────────────────────────


def s27_owasp_table(d):
    s = content_slide(d, "Appendix A · OWASP LLM Top 10 (2025), in full", "A1",
                      "Appendix")
    rows = [
        ("LLM01", "Prompt Injection", "Direct & indirect (planted correspondence)",
         "Core"),
        ("LLM02", "Sensitive Information Disclosure", "Balance/PII exfiltration probes",
         "Core"),
        ("LLM03", "Supply Chain", "Out of scope for this arena", "n/a"),
        ("LLM04", "Data & Model Poisoning", "Poisoned CRM context as attack vector",
         "Touched"),
        ("LLM05", "Improper Output Handling", "Unsafe tool calls from model output",
         "Covered"),
        ("LLM06", "Excessive Agency", "Unauthorised transfers & actions", "Core"),
        ("LLM07", "System Prompt Leakage", "Policy/prompt extraction attempts",
         "Covered"),
        ("LLM08", "Vector & Embedding Weaknesses", "Not applicable to targets", "n/a"),
        ("LLM09", "Misinformation", "Referee flags fabricated assurances", "Touched"),
        ("LLM10", "Unbounded Consumption", "Turn limits guard the loop", "Touched"),
    ]
    y0, rh = 2.0, 0.5
    hdr = ["ID", "Risk", "How Gauntlet exercises it", "Focus"]
    xs = [0.42, 1.7, 5.4, 11.6]
    ws = [1.2, 3.6, 6.1, 1.3]
    text(s, xs[0], y0 - 0.02, ws[0], 0.3, hdr[0], size=11, bold=True, color=WHITE,
         after=0)
    rect(s, 0.42, y0 - 0.08, 12.5, 0.44, fill=INK, radius=0.03)
    for j, htext in enumerate(hdr):
        text(s, xs[j] + 0.1, y0 - 0.02, ws[j], 0.32, htext, size=11, bold=True,
             color=WHITE, after=0)
    for i, (code, name, how, focus) in enumerate(rows):
        yy = y0 + 0.44 + i * rh
        if i % 2 == 0:
            rect(s, 0.42, yy, 12.5, rh, fill=MIST)
        fc = {"Core": ORANGE, "Covered": CORAL, "Touched": TEAL, "n/a": GRAY}[focus]
        text(s, xs[0] + 0.1, yy, ws[0], rh, code, size=11, bold=True, color=INK,
             anchor="m", after=0)
        text(s, xs[1] + 0.1, yy, ws[1], rh, name, size=11, color=INK, anchor="m",
             after=0)
        text(s, xs[2] + 0.1, yy, ws[2], rh, how, size=10.5, color=INKSOFT, anchor="m",
             after=0)
        chip(s, xs[3] + 0.1, yy + 0.09, 1.1, focus, fill=fc, size=9.5, h=0.3)


def s28_mitre_table(d):
    s = content_slide(d, "Appendix B · MITRE ATLAS technique mapping", "A2",
                      "Appendix")
    rows = [
        ("Reconnaissance", "n/a", "Persona probes the agent's role, tools and limits"),
        ("Initial Access", "AML.T0051", "LLM Prompt Injection · direct user channel "
         "attacks"),
        ("Initial Access", "AML.T0051.001", "Indirect injection via planted CRM "
         "correspondence"),
        ("ML Attack Staging", "AML.T0054", "LLM Jailbreak · role play & policy "
         "override personas"),
        ("Exfiltration", "AML.T0057", "LLM Data Leakage · balance & PII disclosure"),
        ("Impact", "n/a", "Unauthorised transfer / action without verification"),
    ]
    y0, rh = 2.15, 0.66
    rect(s, 0.42, y0 - 0.5, 12.5, 0.44, fill=INK, radius=0.03)
    for j, (htext, xx, ww) in enumerate([("ATLAS Tactic", 0.42, 3.0),
                                         ("Technique ID", 3.5, 2.6),
                                         ("What Gauntlet does", 6.2, 6.7)]):
        text(s, xx + 0.1, y0 - 0.48, ww, 0.32, htext, size=11, bold=True, color=WHITE,
             after=0)
    for i, (tactic, tid, desc) in enumerate(rows):
        yy = y0 + i * rh
        if i % 2 == 0:
            rect(s, 0.42, yy, 12.5, rh, fill=MIST)
        text(s, 0.52, yy, 3.0, rh, tactic, size=12, bold=True, color=TEAL, anchor="m",
             after=0)
        idc = GRAY if tid == "n/a" else RUST
        text(s, 3.6, yy, 2.6, rh, tid, size=11.5, bold=True, color=idc, anchor="m",
             font="Consolas", after=0)
        text(s, 6.3, yy, 6.5, rh, desc, size=11.5, color=INKSOFT, anchor="m", after=0)
    footer(s, "IDs reference the public MITRE ATLAS matrix (atlas.mitre.org). Every "
              "recorded fight carries both an OWASP and an ATLAS tag.")


def s29_library(d):
    s = content_slide(d, "Appendix C · the attack library", "A3", "Appendix")
    kpi(s, 0.42, 1.75, 2.7, "8", "red team personas", accent=ORANGE, h=1.4, nsize=38)
    kpi(s, 3.32, 1.75, 2.7, "15", "fight scenarios", accent=TEAL, h=1.4, nsize=38)
    kpi(s, 6.22, 1.75, 2.7, "4", "blue postures", accent=PURPLE, h=1.4, nsize=38)
    kpi(s, 9.12, 1.75, 3.38, "+∞", "invented at runtime", accent=CYAN, h=1.4,
        nsize=38)
    # personas
    text(s, 0.42, 3.45, 6.0, 0.35, "PERSONAS (SAMPLE)", size=12, bold=True,
         color=ORANGE, after=0)
    for i, p in enumerate(["Fake CEO", "Aggressive lawyer", "Distressed customer",
                           "IT support impersonator", "Regulator", "Insider threat"]):
        col, row = i % 2, i // 2
        chip(s, 0.42 + col * 3.05, 3.9 + row * 0.55, 2.85, p, fill=MIST, color=RUST,
             size=11.5, h=0.42, bold=True)
    # blue modes
    text(s, 6.9, 3.45, 6.0, 0.35, "BLUE POSTURES", size=12, bold=True, color=TEAL,
         after=0)
    modes = [("Standard", "strict policy, strong guardrails"),
             ("Lenient", "weaker model, looser policy"),
             ("Naive", "escalation averse, over helpful"),
             ("External", "third party LangGraph, no tools")]
    for i, (m, sub) in enumerate(modes):
        yy = 3.9 + i * 0.62
        rect(s, 6.9, yy, 0.16, 0.44, fill=TEAL)
        text(s, 7.16, yy, 2.0, 0.44, m, size=13, bold=True, color=INK, anchor="m",
             after=0)
        text(s, 9.1, yy, 3.8, 0.44, sub, size=11.5, color=INKSOFT, anchor="m", after=0)
    footer(s, "personas/ (8 YAML) × scenarios/ (15 YAML) × blue modes, "
              "the Coach writes new ones here at runtime.")


def s30_methodology(d):
    s = content_slide(d, "Appendix D · how a verdict is scored", "A4", "Appendix")
    text(s, 0.42, 1.5, 12.4, 0.5,
         "The Referee judges the full transcript plus a tool ledger, what Blue "
         "actually did, not just what it said.", size=15.5, color=INKSOFT, after=0)
    steps = [
        ("Two parallel views", "Red and Blue each see the other as “the user”. "
         "Red opens with a scripted line, then they alternate up to max_turns.", TEAL),
        ("Tool ledger", "Every Blue tool call (transfers, balance disclosure, "
         "fetching planted correspondence) is recorded, not just the chat.",
         ORANGE),
        ("Early stops", "The loop ends on escalation, Red's <<END_CALL>>, or a policy "
         "breach, whichever comes first.", PURPLE),
        ("Verdict", "The Referee scores transcript + ledger into a defense score and a "
         "win/loss, then save_run() writes the JSON.", CYAN),
    ]
    x, w, gap = 0.42, 6.15, 0.2
    for i, (h, b, c) in enumerate(steps):
        col, row = i % 2, i // 2
        lx = x + col * (w + gap)
        ty = 2.4 + row * 1.9
        card(s, lx, ty, w, 1.7, fill=WHITE, line=CLOUD, line_w=1.25, radius=0.04)
        rect(s, lx, ty, 0.12, 1.7, fill=c)
        text(s, lx + 0.35, ty + 0.22, w - 0.6, 0.4, h, size=15, bold=True, color=INK,
             after=0)
        text(s, lx + 0.35, ty + 0.72, w - 0.65, 0.9, b, size=12, color=INKSOFT,
             spacing=1.14, after=0)
    footer(s, "runner.py drives the fight · referee.py judges it · save_run() "
              "persists every fight as JSON in runs/.")


# ─────────────────────────────── build ───────────────────────────────────


# ─────────────────────────── talk track ──────────────────────────────────
# Spoken, per-slide. Order MUST match the build() sequence. Embedded as
# PowerPoint speaker notes and dumped to docs/PITCH-TALK-TRACK.md.

NOTES = [
    ("Title",
     "Hi, I'm Tim. This is Gauntlet, an adversarial test cloud for UiPath agents. "
     "The one line version: we put an attacker in the loop, and turn every attack "
     "that works into a permanent test. Let me show you why that matters."),
    ("Hook",
     "Here's the uncomfortable truth. Your agent passes every test you wrote. It "
     "sails through your evals. Then a real attacker shows up with a prompt you "
     "never imagined, and nothing you built saw it coming. That gap is what "
     "Gauntlet closes."),
    ("Why agents raise the stakes",
     "With agents, that gap is expensive. A chatbot that says something wrong is "
     "embarrassing. An agent that does something wrong is a breach. Agents act, "
     "they move money and read customer records. The attack surface is language "
     "itself, so there's nothing to sanitize. And they're nondeterministic, so one "
     "passing test proves almost nothing."),
    ("The gap today",
     "Today, UiPath Agent Evaluations answer one question well: does my agent pass "
     "the rubric I wrote? That's useful, but it's a mirror of your own blind spots. "
     "No adversary is ever in the loop. Did we test prompt injection comes back as "
     "maybe, not a number. Safety today is hope, not proof."),
    ("Beyond the LLM judge",
     "Now, you might be thinking: Test Cloud already uses an LLM as a judge. True. "
     "So does Gauntlet, that's our Referee. But the judge was never the hard part. "
     "A judge only grades the attacks someone already thought to run. What's "
     "missing is the adversary that invents new ones, self play that raises the bar "
     "every time you pass, and a loop that turns the verdict into a test and a fix. "
     "An LLM judge tells you whether you passed. Gauntlet decides what the exam is."),
    ("The delta",
     "So here's the shift. From static rubrics a human wrote once, to self play "
     "that invents new evals. From maybe we tested that, to an OWASP and MITRE "
     "coverage matrix on every release. From manual repro and manual test cases, to "
     "auto populated Test Manager regressions. And from a developer patching by "
     "hand, to a fix already drafted in Action Center."),
    ("Section: Meet Gauntlet",
     "The idea is simple. If attackers are the problem, we fight fire with fire. "
     "Meet Gauntlet."),
    ("The core loop",
     "One loop. A Red Coach attacks your Blue agent. A Referee scores each round. "
     "When Red wins, that attack becomes a permanent Test Manager regression. When "
     "Blue holds, the Coach invents something harder. And because the Coach is "
     "rewarded for landing attacks, the test suite grows itself. The corpus is the "
     "product."),
    ("Red team vs blue team",
     "This is red team versus blue team, straight out of security. The red team "
     "breaks it: injection, impersonation, exfiltration. The blue team holds the "
     "line: input distrust, identity checks, escalation. Gauntlet automates both "
     "sides and runs them against each other, continuously."),
    ("The self play thesis",
     "Why an AI adversary? Because you cannot hand write every attack. The space is "
     "effectively infinite, and it shifts with every model update. So we do what "
     "AlphaGo did, self play. An adversary that learns and escalates is the only "
     "thing that scales to an AI defender. Safety becomes a number you measure, not "
     "a box you tick."),
    ("Section: How a fight works",
     "Let me walk you through exactly how a fight works."),
    ("The cast",
     "Four agents, four jobs. The Red Coach attacks. The Blue target defends, "
     "that's your agent. The Referee judges. And the Fix Recommender writes the "
     "patch when Blue loses. And it's deliberately polyglot: coded Python agents "
     "and low code Agent Builder agents, fighting in the same ring."),
    ("One round end to end",
     "A single round runs as a Maestro Flow. Red and Blue are matched, they "
     "negotiate over several turns, the Referee delivers a verdict, and the round "
     "is archived. If Red won, it goes to Test Manager. If Blue lost, the Fix "
     "Recommender opens an Action Center task."),
    ("The Coach",
     "The Coach is the clever part. It doesn't attack at random. It uses Thompson "
     "sampling to play the attack most likely to break this particular target. And "
     "when nothing in the corpus lands anymore, it asks the LLM to invent a brand "
     "new persona, conditioned on the last few losses. No human writes the next "
     "test."),
    ("A real attack (the breach)",
     "Here's a real attack it found. The scariest injections never touch the chat "
     "box, they ride in on data your agent is trusted to read. Gauntlet planted a "
     "fake KYC subsystem note in the CRM. The agent fetched it as a tool result, "
     "believed it, and disclosed an $84,320 balance with zero 2FA. That's OWASP "
     "LLM01 arriving as LLM02. Pause on that one."),
    ("The flywheel",
     "And every outcome does work. A win writes a regression into Test Manager. A "
     "loss opens a fix in Action Center, with the patch already drafted. Either "
     "way, the agent gets harder to break and the corpus grows. That's the "
     "flywheel."),
    ("Section: Built on UiPath",
     "And this isn't slideware. It's built and deployed on UiPath."),
    ("Architecture",
     "Here's the whole system. A Coded App console on top. Underneath, a Maestro "
     "Case holds the fight, a Maestro Flow runs each round, with the Coach, your "
     "agent, the Referee, and the Fix Recommender inside it. Outcomes write "
     "straight back to Test Manager and Action Center in your tenant. And the runs "
     "folder, the corpus, is the shared source of truth."),
    ("Framework neutral",
     "It's framework neutral. Low code Agent Builder agents, coded LangGraph "
     "agents, even an external third party target, all fight under the same rules. "
     "One contract, eight UiPath products, zero lock in on how your agent is "
     "built."),
    ("App: Dashboard",
     "Now the actual app. You start in the threat dashboard: robustness score, top "
     "critical findings, a live coverage heatmap, and one click to launch a "
     "fight."),
    ("App: Run a fight",
     "You pick a red persona, a scenario, and the blue posture, standard, lenient, "
     "naive, or an external target, and run a single fight or a batch of up to "
     "ten."),
    ("App: Coach Lab",
     "In Coach Lab you watch the Coach invent an attack live. The LLM call streams "
     "from your browser with your own key, and the risk table shows exactly where "
     "blue is weakest."),
    ("App: Fix Lab",
     "When a fight is lost, Fix Lab already has the answer: root cause, a system "
     "prompt patch, and a regression test. One click files it to Action Center for "
     "human approval."),
    ("App: Audit",
     "And Audit is the payoff. Every fight double tagged to OWASP and MITRE, "
     "rendered as a live coverage matrix, red where blue is weak. The answer a "
     "compliance officer can actually read."),
    ("Section: Coverage you can audit",
     "Which brings me to standards, because coverage you can't audit isn't really "
     "coverage."),
    ("OWASP LLM Top 10",
     "Every fight maps to the OWASP LLM Top 10, the industry standard list of the "
     "worst LLM risks. We focus where agents actually get hit: prompt injection, "
     "sensitive information disclosure, excessive agency, and system prompt "
     "leakage."),
    ("MITRE ATLAS",
     "We tag those same fights to MITRE ATLAS, the ATT&CK of machine learning. "
     "OWASP names the vulnerability, ATLAS names the adversary's behavior. Together "
     "they give a security team the language they already use everywhere else."),
    ("ISO 42001 / NIST / EU AI Act",
     "And this matters for governance. ISO 42001, the NIST AI Risk Management "
     "Framework, and the EU AI Act all expect adversarial testing and records of "
     "it. Gauntlet produces that evidence just by running. It doesn't certify you, "
     "it hands you the artifact an auditor asks for."),
    ("Proof",
     "So does it work? Seventy five fights. Under the strict posture, blue blocked "
     "all 27 attacks, averaging 90.6 out of 100. Weaken the policy, and the "
     "attacker win rate climbs to 22 percent, measured, not guessed. And that "
     "$84,320 breach is one of nine critical findings we auto diagnosed and "
     "proposed a patch for."),
    ("Deployed and live",
     "It's all deployed and live in the UiPath Labs tenant: the solution, the app, "
     "a demo video, and an open MIT licensed repo. The vision is simple. "
     "Continuous adversarial assurance for every agent you ship, the same way you "
     "already run unit tests and CI. An agent that keeps agents honest."),
    ("Close",
     "Go safe, or go home. Thank you. I'd love to take your questions."),
    ("Appendix A: OWASP table",
     "Backup slide. The full OWASP LLM Top 10, with exactly how Gauntlet exercises "
     "each one. Pull it up if a judge wants the detail."),
    ("Appendix B: MITRE techniques",
     "Backup slide. The MITRE ATLAS technique mapping, with the real technique IDs "
     "carried on every fight."),
    ("Appendix C: attack library",
     "Backup slide. The attack library: eight personas, fifteen scenarios, four "
     "blue postures, plus whatever the Coach invents at runtime."),
    ("Appendix D: scoring",
     "Backup slide. How a verdict is actually scored: two parallel views, a tool "
     "ledger of what blue really did, early stop conditions, and the Referee's "
     "final call."),
]


def _apply_notes(d):
    slides = list(d.prs.slides)
    assert len(NOTES) == len(slides), f"NOTES {len(NOTES)} != slides {len(slides)}"
    for slide, (_, script) in zip(slides, NOTES):
        slide.notes_slide.notes_text_frame.text = script


def _write_talk_track():
    head = [
        "# Gauntlet: finalist pitch talk track",
        "",
        "Per slide speaker notes, also embedded in `Gauntlet Pitch Deck.pptx` "
        "(visible in PowerPoint Presenter View).",
        "",
        "**Pacing:** the core deck (slides 1 to 31) runs about 7 to 8 minutes at a "
        "brisk pace. To hit 5 minutes, trim slides 3, 6 and 19. Section dividers "
        "(7, 11, 17, 25) are 3 to 5 second transitions, do not linger. The appendix "
        "(32 to 35) is flip to only.",
        "",
        "**Live demo:** slides 20 to 24 are your natural demo window. If you demo "
        "the app live, drive the app and skip the narration; these notes are your "
        "fallback.",
        "",
        "**Slow down on:** slide 15 (the breach) and slide 29 (the numbers). Let "
        "them land.",
        "",
        "---",
        "",
    ]
    body = []
    for i, (label, script) in enumerate(NOTES, 1):
        body += [f"## Slide {i}: {label}", "", script, ""]
    OUTPUT_MD.write_text("\n".join(head + body))
    print(f"Wrote {OUTPUT_MD}")


def build():
    d = Deck()
    for fn in [s01_title, s02_hook, s03_stakes, s04_gap, s04b_judge, s05_delta,
               s06_section_meet,
               s07_what, s08_redblue, s09_thesis, s10_section_how, s11_cast, s12_round,
               s13_coach, s14_injection, s15_flywheel, s16_section_uipath,
               s17_architecture, s18_components,
               lambda d: app_slide(d, "Start in the Threat Dashboard",
                   "01-dashboard.png", "1", "Dashboard", [
                       "Robustness score and top critical findings at a glance.",
                       "Recent fights and a mini OWASP coverage heatmap.",
                       "Quick action tiles to launch the next fight."]),
               lambda d: app_slide(d, "Launch a fight in two clicks",
                   "03-run-fight.png", "2", "Run a fight", [
                       "Pick a Red persona and a fight scenario.",
                       "Choose the Blue posture: Standard, Lenient, Naive or External.",
                       "Run one fight, or a batch of up to 10."]),
               lambda d: app_slide(d, "Watch the Coach invent an attack, live",
                   "02-coach-lab.png", "3", "Coach Lab", [
                       "The Coach plays the highest expected reward persona, or "
                       "invents a fresh one.",
                       "The LLM call streams from your browser with your own key.",
                       "The risk weighted table shows exactly where Blue is weakest."]),
               lambda d: app_slide(d, "Turn a loss into a filed fix",
                   "04-fix-lab.png", "4", "Fix Lab", [
                       "Open any losing fight and read the drafted patch.",
                       "Root cause, a system prompt patch, and a regression test.",
                       "File it to UiPath Action Center for human approval."]),
               lambda d: app_slide(d, "Prove coverage in the Audit view",
                   "06-audit.png", "5", "Audit", [
                       "Every fight double tagged to OWASP and MITRE ATLAS.",
                       "A live coverage matrix, red where Blue is weak.",
                       "The answer a compliance officer can actually read."]),
               s20_section_audit,
               s21_owasp, s22_mitre, s23_iso, s24_proof, s25_live, s26_close,
               s27_owasp_table, s28_mitre_table, s29_library, s30_methodology]:
        fn(d)
    _apply_notes(d)
    d.save()
    _write_talk_track()
    print(f"Wrote {OUTPUT}  ({len(d.prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
